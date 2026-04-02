#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
生成训练客户端培训 PPT
使用 python-pptx 库生成 PowerPoint 演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os


def create_title_slide(prs, title, subtitle):
    """创建标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide


def create_content_slide(prs, title, content_items):
    """创建内容幻灯片（项目符号列表）"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.level = 0
        if item.startswith("  "):
            p.level = 1
            p.text = item.strip()

    return slide


def create_diagram_slide(prs, title, diagram_text):
    """创建图表幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局
    slide.shapes.title.text = title

    # 添加文本框显示图表
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = diagram_text
    p.font.size = Pt(11)
    p.font.name = 'Courier New'

    return slide


def create_table_slide(prs, title, headers, data):
    """创建表格幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白布局
    slide.shapes.title.text = title

    # 添加表格
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(0.8)

    rows = len(data) + 1
    cols = len(headers)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 设置列宽
    column_widths = [Inches(1.5), Inches(2), Inches(3.5), Inches(2)]
    for i, col_width in enumerate(column_widths[:cols]):
        table.columns[i].width = col_width

    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True

    # 填充数据
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data[:cols]):
            table.cell(row_idx + 1, col_idx).text = str(cell_data)
            table.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(10)

    return slide


def create_code_slide(prs, title, code, description=""):
    """创建代码示例幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    top = Inches(1.5)

    # 添加描述
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.8))
        desc_tf = desc_box.text_frame
        desc_tf.word_wrap = True
        desc_p = desc_tf.paragraphs[0]
        desc_p.text = description
        desc_p.font.size = Pt(14)
        top = Inches(2.3)

    # 添加代码框
    code_box = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(4))
    code_tf = code_box.text_frame
    code_tf.word_wrap = True

    p = code_tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(10)
    p.font.name = 'Courier New'

    return slide


def generate_ppt():
    """生成完整的 PPT"""
    prs = Presentation()

    # 设置幻灯片大小为标准 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("正在创建 PPT...")

    # ========== 幻灯片 1: 标题页 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "训练客户端技术详解"
    slide.placeholders[1].text = "原理 · 架构 · 使用指南\n\n2026 年 3 月"
    print("  - 标题页")

    # ========== 幻灯片 2: 目录 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "目录"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "1. 系统架构概览",
        "2. 消息协议详解",
        "3. TrainingClient 核心功能",
        "4. 完整工作流程",
        "5. 参数配置说明",
        "6. 实战代码示例",
        "7. 常见问题解答"
    ]

    for i, item in enumerate(contents):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.space_after = Pt(12)
    print("  - 目录页")

    # ========== 幻灯片 3: 系统架构概览 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "1. 系统架构概览"

    # 绘制架构图
    shapes = slide.shapes

    # 客户端框
    left = Inches(1)
    top = Inches(2)
    width = Inches(3)
    height = Inches(1.5)
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.text = "训练客户端\n(training_client.py)"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(79, 129, 189)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 网络框
    left = Inches(5)
    shape = shapes.add_shape(MSO_SHAPE.CLOUD, left, top, width, height)
    shape.text = "TCP/IP 网络\n127.0.0.1:8888"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(146, 208, 80)

    # 服务器框
    left = Inches(9)
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.text = "训练服务器\n(model_training.py)"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(237, 125, 49)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 箭头
    top_arrow = Inches(3.2)
    shape = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.2), top_arrow, Inches(0.8), Inches(0.3))
    shape.text_frame.paragraphs[0].text = "请求"
    shape.text_frame.paragraphs[0].font.size = Pt(10)

    shape = shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(8), top_arrow, Inches(0.8), Inches(0.3))
    shape.text_frame.paragraphs[0].text = "响应"
    shape.text_frame.paragraphs[0].font.size = Pt(10)

    print("  - 系统架构概览")

    # ========== 幻灯片 4: 消息协议结构 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "2. 消息协议结构"

    # 协议结构图
    top = Inches(2)

    # 包头
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(2.5), Inches(1))
    shape.text = "包头 (4 字节)\n0x55AA55AA"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(68, 114, 196)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 长度
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), top, Inches(2.5), Inches(1))
    shape.text = "长度 (4 字节)\n数据字节数"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(146, 208, 80)

    # 校验和
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), top, Inches(2.5), Inches(1))
    shape.text = "校验和 (4 字节)\n完整性验证"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(237, 125, 49)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 数据
    shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), top, Inches(4.5), Inches(1))
    shape.text = "数据 (N 字节)\nJSON 字符串"
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(91, 155, 213)
    shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # 说明
    tf = shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(3)).text_frame
    tf.word_wrap = True

    protocol_desc = """
协议详解：

1. 包头 (Header): 固定值 0x55AA55AA，用于识别消息起始位置
   - 大端序 4 字节无符号整数
   - 接收方通过此值判断消息边界

2. 长度 (Length): 后续数据部分的字节数
   - 大端序 4 字节无符号整数
   - 用于确定需要读取的字节数

3. 校验和 (Checksum): 数据部分的累加和
   - 计算方法：所有字节相加，取低 32 位
   - 用于验证数据完整性

4. 数据 (Data): UTF-8 编码的 JSON 字符串
   - 包含命令类型、请求 ID、参数等
"""
    p = tf.paragraphs[0]
    p.text = protocol_desc
    p.font.size = Pt(12)
    print("  - 消息协议结构")

    # ========== 幻灯片 5: MessageProtocol 类方法 ==========
    headers = ["方法", "参数", "返回值", "说明"]
    data = [
        ["calculate_checksum", "data: bytes", "int", "计算数据的校验和"],
        ["pack", "message: str", "bytes", "将消息打包成二进制"],
        ["receive", "conn, timeout", "str | None", "从连接接收并解析消息"]
    ]
    slide = create_table_slide(prs, "MessageProtocol 类方法", headers, data)
    print("  - MessageProtocol 方法表")

    # ========== 幻灯片 6: 消息打包示例 ==========
    code = '''
# 消息打包过程
message = {"command": "ping", "request_id": "123"}
packed = MessageProtocol.pack(message)

# 打包后的二进制结构：
# [0x55AA55AA][长度：44][校验和：0x1234ABCD][{"command":...}]
#    4 字节      4 字节     4 字节          数据部分

# 校验和计算原理：
data = b'{"command": "ping"}'
checksum = 0
for byte in data:
    checksum = (checksum + byte) & 0xFFFFFFFF  # 保持 32 位
'''
    slide = create_code_slide(prs, "消息打包代码示例", code, "理解 pack() 方法的工作原理")
    print("  - 消息打包示例")

    # ========== 幻灯片 7: TrainingClient 类结构 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. TrainingClient 核心属性"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    attrs = [
        "self.host - 服务器地址（默认：127.0.0.1）",
        "self.port - 服务器端口（默认：8888）",
        "self.socket - TCP Socket 对象",
        "self.connected - 连接状态（bool）",
        "self.request_counter - 请求计数器（用于生成唯一 ID）",
        "self.callbacks - 回调函数字典"
    ]

    for i, item in enumerate(attrs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.space_after = Pt(10)
    print("  - TrainingClient 属性")

    # ========== 幻灯片 8: 核心方法（一） ==========
    headers = ["方法", "参数", "返回值", "作用"]
    data = [
        ["connect()", "无", "bool", "连接到训练服务器"],
        ["disconnect()", "无", "无", "断开与服务器的连接"],
        ["ping()", "无", "bool", "测试连接是否通畅"],
        ["send_request()", "command, data", "dict | None", "发送请求并获取响应"]
    ]
    slide = create_table_slide(prs, "核心方法（连接类）", headers, data)
    print("  - 核心方法表 1")

    # ========== 幻灯片 9: 核心方法（二） ==========
    headers = ["方法", "参数", "返回值", "作用"]
    data = [
        ["create_task()", "params: dict", "task_id | None", "创建新的训练任务"],
        ["start_training()", "task_id: str", "bool", "启动指定的训练任务"],
        ["stop_training()", "task_id: str", "bool", "停止正在运行的训练"],
        ["get_task_status()", "task_id: str", "dict | None", "获取任务当前状态"]
    ]
    slide = create_table_slide(prs, "核心方法（任务管理）", headers, data)
    print("  - 核心方法表 2")

    # ========== 幻灯片 10: 核心方法（三） ==========
    headers = ["方法", "参数", "返回值", "作用"]
    data = [
        ["get_progress()", "task_id: str", "dict | None", "获取训练进度详情"],
        ["long_poll_progress()", "task_id, timeout", "dict | None", "长轮询等待进度更新"],
        ["list_tasks()", "无", "dict | None", "列出所有训练任务"],
        ["delete_task()", "task_id: str", "bool", "删除指定的训练任务"],
        ["monitor_training()", "task_id, interval", "无", "持续监控训练进度"]
    ]
    slide = create_table_slide(prs, "核心方法（监控类）", headers, data)
    print("  - 核心方法表 3")

    # ========== 幻灯片 11: 完整工作流程图 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "4. 完整工作流程图"

    tf = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)).text_frame
    tf.word_wrap = True

    flow_chart = """
┌─────────────────────────────────────────────────────────────────────────┐
│                            训练客户端工作流程                            │
└─────────────────────────────────────────────────────────────────────────┘

  ┌──────────────┐     ┌──────────────┐     ┌──────────────┐
  │   1. 创建    │────>│   2. 连接    │────>│   3. 验证    │
  │ TrainingClient│     │   connect()  │     │    ping()    │
  └──────────────┘     └──────────────┘     └──────────────┘
                                                 │
                                                 ▼
  ┌──────────────┐     ┌──────────────┐     ┌──────────────┐
  │   6. 断开    │<────│   5. 监控    │<────│   4. 创建    │
  │ disconnect() │     │ monitor()    │     │  create_task │
  └──────────────┘     └──────────────┘     └──────────────┘
                              │                     │
                              ▼                     ▼
                       ┌──────────────┐     ┌──────────────┐
                       │  get_progress│     │ start_training│
                       └──────────────┘     └──────────────┘


  详细步骤说明：

  1️⃣  创建：实例化 TrainingClient，指定服务器地址和端口
  2️⃣  连接：调用 connect() 建立 TCP 连接
  3️⃣  验证：调用 ping() 确认服务器响应正常
  4️⃣  创建任务：调用 create_task(params) 配置训练参数
  5️⃣  启动训练：调用 start_training(task_id) 开始训练
  6️⃣  监控进度：循环调用 get_progress() 或 monitor_training()
  7️⃣  停止训练：调用 stop_training(task_id)（可选，训练完成可跳过）
  8️⃣  断开连接：调用 disconnect() 释放资源
"""
    p = tf.paragraphs[0]
    p.text = flow_chart
    p.font.size = Pt(9)
    p.font.name = 'Courier New'
    print("  - 工作流程图")

    # ========== 幻灯片 12: create_task 参数详解 ==========
    headers = ["参数", "类型", "默认值", "说明"]
    data = [
        ["model_type", "str", "'detect'", "模型类型：detect/segment/classify"],
        ["image_size", "int", "640", "输入图像尺寸（像素）"],
        ["dataset", "str", "'data'", "训练数据集路径"],
        ["epochs", "int", "50", "训练轮次数"],
        ["batch_size", "int", "32", "每批次样本数"],
        ["learning_rate", "float", "0.001", "学习率"],
        ["trainset_ratio", "float", "0.9", "训练集占总数据的比例"]
    ]
    slide = create_table_slide(prs, "5. create_task 参数详解", headers, data)
    print("  - create_task 参数")

    # ========== 幻灯片 13: 参数影响说明 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "参数对训练的影响"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    params_effect = [
        "epochs（训练轮次）:",
        "  • 太少：模型可能未充分学习（欠拟合）",
        "  • 太多：可能过拟合，训练时间增长",
        "  • 建议：根据数据集大小调整，一般 30-100",
        "",
        "batch_size（批次大小）:",
        "  • 太小：训练不稳定，但内存占用少",
        "  • 太大：内存需求高，收敛速度可能变慢",
        "  • 建议：16/32/64，根据 GPU 显存调整",
        "",
        "learning_rate（学习率）:",
        "  • 太小：收敛速度慢",
        "  • 太大：可能不收敛或震荡",
        "  • 建议：0.001-0.01，可使用学习率衰减"
    ]

    for i, line in enumerate(params_effect):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16) if not line.startswith("  ") else Pt(14)
        p.space_after = Pt(8)
    print("  - 参数影响说明")

    # ========== 幻灯片 14: 使用示例 - 基础用法 ==========
    code = '''
from training_client import TrainingClient

# 1. 创建客户端
client = TrainingClient(host="127.0.0.1", port=8888)

# 2. 连接到服务器
if not client.connect():
    print("连接失败")
    exit()

# 3. 测试连接
if not client.ping():
    print("服务器无响应")
    exit()

# 4. 创建训练任务
params = {
    "model_type": "detect",
    "image_size": 640,
    "dataset": "data/my_dataset",
    "epochs": 50,
    "batch_size": 32,
    "learning_rate": 0.001,
    "trainset_ratio": 0.9
}
task_id = client.create_task(params)

# 5. 启动训练
client.start_training(task_id)

# 6. 监控进度
while True:
    progress = client.get_progress(task_id)
    if progress:
        epoch = progress["epoch"]
        total = progress["total_epochs"]
        print(f"进度：{epoch}/{total}")
        if epoch >= total:
            break
    time.sleep(2)

# 7. 断开连接
client.disconnect()
'''
    slide = create_code_slide(prs, "6. 基础使用示例", code)
    print("  - 基础使用示例")

    # ========== 幻灯片 15: 使用示例 - 监控训练 ==========
    code = '''
def monitor_with_callback(task_id):
    """带回调的监控函数"""

    def on_epoch_update(data):
        """每轮更新时的回调"""
        print(f"轮次 {data['epoch']}/{data['total_epochs']}")
        print(f"  损失：{data['loss']:.4f}")
        print(f"  准确率：{data['accuracy']:.4f}")

    print("开始监控训练...")

    while True:
        # 获取任务状态
        status = client.get_task_status(task_id)
        print(f"状态：{status['status']}")

        # 获取详细进度
        progress = client.get_progress(task_id)
        if progress and "progress" in progress:
            p = progress["progress"]
            on_epoch_update(p)

            # 检查是否完成
            if p.get("epoch", 0) >= p.get("total_epochs", 1):
                print("训练完成!")
                break

        time.sleep(2)  # 每 2 秒查询一次
'''
    slide = create_code_slide(prs, "监控训练进度示例", code)
    print("  - 监控示例")

    # ========== 幻灯片 16: 使用示例 - 错误处理 ==========
    code = '''
def train_with_error_handling():
    """带错误处理的完整示例"""
    client = TrainingClient()

    try:
        # 连接
        if not client.connect():
            print("❌ 连接失败")
            return False

        # 创建任务
        task_id = client.create_task({
            "model_type": "detect",
            "epochs": 50
        })

        if not task_id:
            print("❌ 创建任务失败")
            return False

        # 启动训练
        if not client.start_training(task_id):
            print("❌ 启动失败")
            return False

        # 监控（带超时）
        timeout = 3600  # 1 小时超时
        start_time = time.time()

        while time.time() - start_time < timeout:
            status = client.get_task_status(task_id)
            if status["status"] == "completed":
                print("✅ 训练完成")
                return True
            elif status["status"] == "failed":
                print("❌ 训练失败")
                return False
            time.sleep(5)

        print("⚠️ 超时")
        return False

    except Exception as e:
        print(f"❌ 异常：{e}")
        return False
    finally:
        client.disconnect()
'''
    slide = create_code_slide(prs, "错误处理示例", code)
    print("  - 错误处理示例")

    # ========== 幻灯片 17: 消息传递原理 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "消息传递原理"

    tf = shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5)).text_frame
    tf.word_wrap = True

    msg_flow = """
┌─────────────────────────────────────────────────────────────────┐
│                      客户端 → 服务器 请求                        │
├─────────────────────────────────────────────────────────────────┤
│                                                                  │
│  原始请求: {"command": "create_task", "params": {...}}          │
│              ↓ json.dumps()                                      │
│  JSON 字符串："{"command": "create_task", ...}"                   │
│              ↓ MessageProtocol.pack()                            │
│  二进制：[Header][Length][Checksum][Data]                        │
│              ↓ socket.sendall()                                  │
│  网络传输：TCP 数据流                                              │
│                                                                  │
└─────────────────────────────────────────────────────────────────┘

┌─────────────────────────────────────────────────────────────────┐
│                      服务器 → 客户端 响应                        │
├─────────────────────────────────────────────────────────────────┤
│                                                                  │
│  网络接收：TCP 数据流                                             │
│              ↓ MessageProtocol.receive()                         │
│  二进制：[Header][Length][Checksum][Data]                        │
│              ↓ 验证包头、校验和                                   │
│  JSON 字符串："{"code": 100, "data": {...}}"                      │
│              ↓ json.loads()                                      │
│  Python 对象：{"code": 100, "data": {...}}                        │
│                                                                  │
└─────────────────────────────────────────────────────────────────┘
"""
    p = tf.paragraphs[0]
    p.text = msg_flow
    p.font.size = Pt(10)
    p.font.name = 'Courier New'
    print("  - 消息传递原理")

    # ========== 幻灯片 18: 响应码说明 ==========
    headers = ["响应码", "含义", "场景"]
    data = [
        ["100", "成功", "请求正常处理"],
        ["200", "部分成功", "操作部分完成"],
        ["400", "请求错误", "参数无效或缺失"],
        ["404", "未找到", "任务 ID 不存在"],
        ["500", "服务器错误", "内部异常"],
        ["503", "服务不可用", "服务器过载或维护"]
    ]
    slide = create_table_slide(prs, "常见响应码", headers, data)
    print("  - 响应码表")

    # ========== 幻灯片 19: InteractiveClient ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "InteractiveClient 交互式客户端"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    interactive_desc = """
InteractiveClient 是 TrainingClient 的封装，提供命令行交互界面：

主要功能:
  • 菜单驱动的用户界面
  • 交互式参数输入
  • 自动任务管理（记录当前任务 ID）
  • 友好的错误提示

使用方式:
  python training_client.py --interactive

  或直接运行:
  python training_client.py

菜单选项:
  1. 创建训练任务
  2. 开始训练
  3. 监控训练进度
  4. 停止训练
  5. 查看任务状态
  6. 列出所有任务
  7. 删除任务
  8. 退出
"""
    p = tf.paragraphs[0]
    p.text = interactive_desc
    p.font.size = Pt(16)
    print("  - InteractiveClient")

    # ========== 幻灯片 20: 常见问题 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "7. 常见问题解答"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    faqs = [
        "Q1: 连接失败怎么办？",
        "  A: 检查服务器是否启动，确认地址和端口正确",
        "",
        "Q2: 任务创建失败？",
        "  A: 检查参数格式，确保 dataset 路径存在",
        "",
        "Q3: 训练进度一直为 0？",
        "  A: 确认任务已启动（调用 start_training）",
        "",
        "Q4: 如何停止训练？",
        "  A: 调用 stop_training(task_id) 或直接删除任务",
        "",
        "Q5: 长轮询和普通查询的区别？",
        "  A: 长轮询会等待数据更新，减少无效请求"
    ]

    for i, line in enumerate(faqs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18) if line.startswith("Q") else Pt(14)
        p.space_after = Pt(10)
    print("  - 常见问题")

    # ========== 幻灯片 21: 总结 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "总结"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    summary = [
        "核心知识点:",
        "",
        "1. 消息协议：包头 + 长度 + 校验和 + 数据",
        "2. TCP 通信：基于 Socket 的长连接",
        "3. 请求 - 响应模式：每个请求都有唯一 ID",
        "4. 任务管理：创建 → 启动 → 监控 → 完成",
        "",
        "关键类:",
        "",
        "• MessageProtocol - 消息编解码",
        "• TrainingClient - 核心客户端",
        "• InteractiveClient - 交互式界面",
        "",
        "下一步:",
        "• 阅读 model_training.py 了解服务器端实现",
        "• 尝试修改参数观察训练效果"
    ]

    for i, line in enumerate(summary):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18) if not line.startswith("•") and not line.startswith("  ") else Pt(14)
        p.space_after = Pt(8)
    print("  - 总结页")

    # 保存 PPT
    output_file = "训练客户端技术详解.pptx"
    prs.save(output_file)
    print(f"\nPPT 已保存：{output_file}")
    return output_file


def generate_beginner_ppt():
    """生成小白友好的 PPT 教程"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    print("正在创建小白教程 PPT...")

    # ===== 第 1 页：封面 =====
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TrainingClient 使用教程"
    slide.placeholders[1].text = "从零开始，30 分钟上手\nAI 模型训练客户端详解\n\n2026 年 3 月"
    print("  - 封面")

    # ===== 第 2 页：这是什么 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. 这是什么？能做什么？"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "TrainingClient 是一个 Python 类，用来远程控制 AI 训练服务器",
        "",
        "想象一下：你有一个'遥控器'🎮",
        "  • 告诉它训练什么模型",
        "  • 启动/停止训练",
        "  • 查看训练进度",
        "",
        "这个'遥控器'就是 TrainingClient！",
        "",
        "适用场景：",
        "  • 在 Labelme 中标注完图片后，一键训练",
        "  • 批量训练多个模型",
        "  • 远程监控训练进度"
    ]

    for i, item in enumerate(contents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20) if not item.startswith("  ") else Pt(16)
        p.space_after = Pt(10)
    print("  - 这是什么")

    # ===== 第 3 页：工作原理 =====
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "2. 工作原理：其实很简单！"

    # 绘制简单架构图
    client = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(2.5), Inches(1.2))
    client.fill.solid()
    client.fill.fore_color.rgb = RGBColor(222, 235, 247)
    client.line.color.rgb = RGBColor(79, 129, 189)
    tf_client = client.text_frame
    tf_client.text = "你的 Python 代码\n(TrainingClient)"
    tf_client.paragraphs[0].font.size = Pt(12)
    tf_client.paragraphs[0].alignment = PP_ALIGN.CENTER

    server = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2), Inches(2.5), Inches(1.2))
    server.fill.solid()
    server.fill.fore_color.rgb = RGBColor(226, 240, 224)
    server.line.color.rgb = RGBColor(92, 148, 69)
    tf_server = server.text_frame
    tf_server.text = "训练服务器\n(model_training.py)"
    tf_server.paragraphs[0].font.size = Pt(12)
    tf_server.paragraphs[0].alignment = PP_ALIGN.CENTER

    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.5), Inches(2.3), Inches(3), Inches(0.25))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = RGBColor(255, 200, 100)
    arrow1.text_frame.text = "发送命令 (JSON)"
    arrow1.text_frame.paragraphs[0].font.size = Pt(10)
    arrow1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    arrow2 = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.5), Inches(2.7), Inches(3), Inches(0.25))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(180, 220, 180)
    arrow2.text_frame.text = "返回结果"
    arrow2.text_frame.paragraphs[0].font.size = Pt(10)
    arrow2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    net = slide.shapes.add_textbox(Inches(4.2), Inches(3.1), Inches(1.6), Inches(0.5))
    net.text_frame.text = "TCP/IP 网络\n127.0.0.1:8888"
    net.text_frame.paragraphs[0].font.size = Pt(11)
    net.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 说明
    note = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
    tf = note.text_frame
    tf.text = "通信过程（3 步）："
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True

    steps = [
        "1. 你的代码发送 JSON 命令（如'开始训练'）",
        "2. 服务器执行命令，返回结果",
        "3. 你的代码解析结果，显示给用户",
        "",
        "好消息：这些细节已经封装好了，你只需要调用方法就行！"
    ]
    for i, text in enumerate(steps):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(14)
        p.space_before = Pt(6)
        if "好消息" in text:
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 128, 0)
    print("  - 工作原理")

    # ===== 第 4 页：核心功能 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. 核心功能一览"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "TrainingClient 能做什么？（10 个方法）",
        "",
        "连接相关：",
        "  connect()        - 连接到训练服务器",
        "  disconnect()     - 断开连接",
        "  ping()           - 测试连接是否正常",
        "",
        "任务管理：",
        "  create_task()    - 创建训练任务",
        "  start_training() - 开始训练",
        "  stop_training()  - 停止训练",
        "  list_tasks()     - 查看所有任务",
        "  delete_task()    - 删除任务",
        "",
        "进度监控：",
        "  get_progress()   - 获取训练进度",
        "  monitor_training() - 持续监控（自动刷新）"
    ]

    for i, item in enumerate(contents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    print("  - 核心功能")

    # ===== 第 5 页：快速开始 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4. 快速开始：3 步上手"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "第 1 步：导入并创建客户端",
        '  from training_client import TrainingClient',
        '  client = TrainingClient()',
        "",
        "第 2 步：连接服务器",
        '  if not client.connect():',
        '      print("连接失败")',
        "",
        "第 3 步：使用功能",
        '  client.ping()  # 测试连接',
        '  task_id = client.create_task({...})  # 创建任务',
        '  client.start_training(task_id)  # 开始训练',
        "",
        "就是这么简单！"
    ]

    for i, item in enumerate(contents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if not item.startswith("  ") else Pt(14)
        p.font.name = "Courier New" if item.startswith("  ") else "Microsoft YaHei"
        p.space_after = Pt(8)
    print("  - 快速开始")

    # ===== 第 6 页：完整示例 =====
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "5. 完整使用示例"

    code = '''# 1. 导入
from training_client import TrainingClient

# 2. 创建客户端
client = TrainingClient(host="127.0.0.1", port=8888)

# 3. 连接
if client.connect():
    print("连接成功！")

    # 4. 测试连接
    if client.ping():
        print("服务器正常")

    # 5. 创建训练任务
    params = {
        "model_type": "detect",      # 模型类型
        "epochs": 50,                 # 训练轮次
        "batch_size": 32,             # 批次大小
        "learning_rate": 0.001        # 学习率
    }
    task_id = client.create_task(params)
    print(f"任务 ID: {task_id}")

    # 6. 开始训练
    client.start_training(task_id)

    # 7. 监控进度（自动刷新）
    client.monitor_training(task_id, poll_interval=2)

    # 8. 断开
    client.disconnect()'''

    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))
    tf = code_box.text_frame
    tf.text = code
    for p in tf.paragraphs:
        p.font.name = "Courier New"
        p.font.size = Pt(11)

    # 右侧说明
    note = slide.shapes.add_textbox(Inches(6.7), Inches(1.3), Inches(2.8), Inches(5.5))
    tf = note.text_frame
    tf.text = "代码说明："
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True

    notes = [
        "• 创建客户端：指定服务器地址和端口",
        "• connect()：建立 TCP 连接",
        "• ping()：发送测试请求",
        "• create_task()：创建任务，返回 task_id",
        "• start_training()：开始训练",
        "• monitor_training()：每 2 秒自动刷新进度",
        "• disconnect()：释放资源"
    ]

    for i, text in enumerate(notes):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(11)
        p.space_before = Pt(6)
    print("  - 完整示例")

    # ===== 第 7 页：InteractiveClient =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "6. 更简单：使用 InteractiveClient"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "不想自己写代码？使用交互式客户端！",
        "",
        "一行代码搞定：",
        '  from training_client import InteractiveClient',
        '  InteractiveClient().run()',
        "",
        "然后会出现菜单：",
        "  1. 创建训练任务",
        "  2. 开始训练",
        "  3. 监控训练进度",
        "  4. 停止训练",
        "  5. 查看任务状态",
        "  6. 列出所有任务",
        "  7. 删除任务",
        "",
        "命令行运行：python training_client.py --interactive"
    ]

    for i, item in enumerate(contents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18) if "不想" in item or "一行" in item or "命令行" in item else Pt(15)
        p.space_after = Pt(8)
    print("  - InteractiveClient")

    # ===== 第 8 页：参数详解 =====
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "7. create_task 参数详解"

    headers = ["参数", "类型", "默认值", "说明"]
    data = [
        ["model_type", "str", "detect", "模型类型：detect/segment/classify"],
        ["image_size", "int", "640", "输入图像尺寸（像素）"],
        ["dataset", "str", "data", "训练数据集路径"],
        ["epochs", "int", "50", "训练轮次数"],
        ["batch_size", "int", "32", "每批次样本数"],
        ["learning_rate", "float", "0.001", "学习率"],
        ["trainset_ratio", "float", "0.9", "训练集占比"]
    ]

    # 创建表格
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(8.8)
    height = Inches(0.8)

    rows = len(data) + 1
    cols = len(headers)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 设置列宽
    column_widths = [Inches(1.4), Inches(0.9), Inches(1.2), Inches(5.3)]
    for i, col_width in enumerate(column_widths[:cols]):
        table.columns[i].width = col_width

    # 设置表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True

    # 填充数据
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data[:cols]):
            table.cell(row_idx + 1, col_idx).text = str(cell_data)
            table.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(11)

    print("  - 参数详解")

    # ===== 第 9 页：参数影响 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "8. 参数对训练的影响"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "epochs（训练轮次）:",
        "  • 太少：模型可能未充分学习（欠拟合）",
        "  • 太多：可能过拟合，训练时间增长",
        "  • 建议：根据数据集大小调整，一般 30-100",
        "",
        "batch_size（批次大小）:",
        "  • 太小：训练不稳定，但内存占用少",
        "  • 太大：内存需求高，收敛速度可能变慢",
        "  • 建议：16/32/64，根据 GPU 显存调整",
        "",
        "learning_rate（学习率）:",
        "  • 太小：收敛速度慢",
        "  • 太大：可能不收敛或震荡",
        "  • 建议：0.001-0.01"
    ]

    for i, item in enumerate(contents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if item.endswith(":") else Pt(14)
        p.space_after = Pt(8)
    print("  - 参数影响")

    # ===== 第 10 页：监控进度 =====
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "9. 如何监控训练进度？"

    code = '''# 方法 1：自动监控（推荐新手）
client.monitor_training(task_id, poll_interval=2)
# 每 2 秒自动刷新一次，显示轮次、损失、准确率

# 方法 2：手动获取
status = client.get_task_status(task_id)
print(f"状态：{status['status']}")
print(f"进度：{status['progress']*100:.1f}%")

# 方法 3：长轮询（节省资源）
progress = client.long_poll_progress(task_id, timeout=30)
# 服务器有更新时才返回，不用频繁查询'''

    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3))
    tf = code_box.text_frame
    tf.text = code
    for p in tf.paragraphs:
        p.font.name = "Courier New"
        p.font.size = Pt(12)

    # 输出示例
    output_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
    tf = output_box.text_frame
    tf.text = "监控时会看到这样的输出："
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True

    output_example = '''
[14:30:25] 轮次 1/50 | 损失：0.5234 | 准确率：0.7123
[14:30:45] 轮次 2/50 | 损失：0.4521 | 准确率：0.7856
[14:31:05] 轮次 3/50 | 损失：0.3892 | 准确率：0.8234
...
[16:45:30] 轮次 50/50 | 损失：0.0823 | 准确率：0.9654
训练完成！'''

    p = tf.add_paragraph()
    p.text = output_example
    p.font.name = "Courier New"
    p.font.size = Pt(10)
    print("  - 监控进度")

    # ===== 第 11 页：错误处理 =====
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "10. 错误处理（很重要！）"

    code = '''# 连接失败处理
client = TrainingClient()
if not client.connect():
    print("连接失败，请检查服务器是否运行")
    return

# 命令失败处理
task_id = client.create_task(params)
if task_id is None:
    print("创建任务失败，检查参数")
    return

# 完整示例（带异常处理）
try:
    client = TrainingClient()
    if not client.connect():
        return
    task_id = client.create_task(params)
    if not task_id:
        return
    client.start_training(task_id)
except Exception as e:
    print(f"出错了：{e}")
finally:
    client.disconnect()  # 确保断开连接'''

    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.5), Inches(5))
    tf = code_box.text_frame
    tf.text = code
    for p in tf.paragraphs:
        p.font.name = "Courier New"
        p.font.size = Pt(11)

    # 提示
    tip_box = slide.shapes.add_textbox(Inches(7.2), Inches(1.3), Inches(2.5), Inches(5))
    tf = tip_box.text_frame
    tf.text = "重要提示："
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True

    tips = [
        "1. 始终检查返回值",
        "2. 使用 try-except 捕获异常",
        "3. 在 finally 中调用 disconnect()",
        "4. 连接失败先检查服务器",
        "5. 任务失败查看错误信息"
    ]

    for i, text in enumerate(tips):
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(12)
        p.space_before = Pt(8)
    print("  - 错误处理")

    # ===== 第 12 页：响应码 =====
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "11. 响应码说明"

    headers = ["响应码", "含义", "怎么办"]
    data = [
        ["100", "成功", "继续执行"],
        ["200", "部分成功", "检查返回数据"],
        ["400", "参数错误", "检查输入参数"],
        ["404", "任务不存在", "确认 task_id"],
        ["500", "服务器错误", "查看服务器日志"]
    ]

    table = slide.shapes.add_table(6, 3, Inches(0.5), Inches(1.5), Inches(9), Inches(1)).table
    column_widths = [Inches(1.5), Inches(2.5), Inches(5)]
    for i, col_width in enumerate(column_widths):
        table.columns[i].width = col_width

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            table.cell(row_idx + 1, col_idx).text = str(cell_data)
            table.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(12)

    # 判断示例
    code = '''# 如何判断是否成功：
response = client.send_request("create_task", {"params": params})
if response.get("code") == 100:
    print("操作成功")
    task_id = response.get("data", {}).get("task_id")
else:
    error_msg = response.get("message", "未知错误")
    print(f"失败：{error_msg}")'''

    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(2.5))
    tf = code_box.text_frame
    tf.text = code
    for p in tf.paragraphs:
        p.font.name = "Courier New"
        p.font.size = Pt(11)
    print("  - 响应码")

    # ===== 第 13 页：实战场景 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "12. 实战应用场景"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "场景 1：在 Labelme 中标注后一键训练",
        "  # 标注完成后",
        '  client = TrainingClient()',
        '  client.connect()',
        '  task_id = client.create_task({"epochs": 50})',
        '  client.start_training(task_id)',
        "",
        "场景 2：批量训练多个模型",
        "  for model_type in ['detect', 'segment']:",
        '      task_id = client.create_task({"model_type": model_type})',
        '      client.start_training(task_id)',
        "",
        "场景 3：自动化测试",
        "  assert client.ping() == True",
        "  assert task_id is not None"
    ]

    for i, item in enumerate(contents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.name = "Courier New" if item.startswith("  ") else "Microsoft YaHei"
        p.space_after = Pt(6)
    print("  - 实战场景")

    # ===== 第 14 页：调试技巧 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "13. 调试技巧"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "1. 确保服务器在运行",
        "   命令行运行：python model_training.py",
        "",
        "2. 先用 ping() 测试连接",
        '   if not client.ping():',
        '       print("服务器无响应")',
        "",
        "3. 查看控制台日志",
        "   代码会打印详细的连接和执行信息",
        "",
        "4. 使用 InteractiveClient 测试",
        "   交互式界面更容易观察问题",
        "",
        "5. 检查端口是否被占用",
        "   Windows: netstat -ano | findstr 8888"
    ]

    for i, item in enumerate(contents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16) if (len(item) > 0 and item[0].isdigit()) else Pt(14)
        p.space_after = Pt(8)
    print("  - 调试技巧")

    # ===== 第 15 页：常见问题 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "14. 常见问题 FAQ"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "Q: 连接失败怎么办？",
        "A: 检查服务器是否运行，端口 8888 是否被占用",
        "",
        "Q: 任务创建失败？",
        "A: 检查 params 参数是否正确，数据集路径是否存在",
        "",
        "Q: 训练进度一直是 0%？",
        "A: 确认调用了 start_training()，等待几秒再查询",
        "",
        "Q: 如何停止训练？",
        "A: 调用 stop_training(task_id) 或直接删除任务",
        "",
        "Q: 训练需要多长时间？",
        "A: 根据数据集大小和 epochs，一般 10 分钟到几小时"
    ]

    for i, item in enumerate(contents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18) if item.startswith("Q:") else Pt(15)
        p.space_after = Pt(10)
    print("  - 常见问题")

    # ===== 第 16 页：总结 =====
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "15. 总结"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    contents = [
        "核心要点：",
        "",
        "  1. TrainingClient 是远程控制训练服务器的'遥控器'",
        "  2. 3 步上手：创建 → 连接 → 使用",
        "  3. 新手推荐用 InteractiveClient，有菜单提示",
        "  4. 记得检查返回值，处理错误",
        "",
        "下一步：",
        "",
        "  • 运行 python training_client.py 体验",
        "  • 参考测试代码：test_training_client.py",
        "  • 查看服务器代码：model_training.py",
        "",
        "祝你学习愉快！🚀"
    ]

    for i, item in enumerate(contents):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20) if "核心" in item or "下一步" in item else Pt(16)
        p.space_after = Pt(10)
    print("  - 总结")

    # 保存
    output_file = "TrainingClient 使用教程_小白版.pptx"
    prs.save(output_file)
    print(f"\nPPT 已保存：{output_file}")
    print(f"共 {len(prs.slides)} 页")
    return output_file


if __name__ == "__main__":
    # 生成小白教程 PPT
    ppt_file = generate_beginner_ppt()
    print(f"\n生成完成：{ppt_file}")
